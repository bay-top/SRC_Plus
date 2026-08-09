#!/usr/bin/env python3
from __future__ import annotations

import argparse
import json
from pathlib import Path

from PIL import Image, ImageDraw
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt
from editorial_rules import load_editorial_rules, validate_manifest

WHITE = RGBColor(255, 255, 255)


def fit_cover(source: Path, target: Path, width: int = 1350, height: int = 1800) -> None:
    image = Image.open(source).convert("RGB")
    ratio = max(width / image.width, height / image.height)
    resized = image.resize((round(image.width * ratio), round(image.height * ratio)), Image.Resampling.LANCZOS)
    left = (resized.width - width) // 2
    top = (resized.height - height) // 2
    resized.crop((left, top, left + width, top + height)).save(target, quality=95)


def gradient_overlay(base: Image.Image, top_end: float, bottom_start: float, top_alpha: int, bottom_alpha: int) -> Image.Image:
    rgba = base.convert("RGBA")
    overlay = Image.new("RGBA", rgba.size, (0, 0, 0, 0))
    draw = ImageDraw.Draw(overlay)
    h = rgba.height
    top_px = max(1, int(h * top_end))
    for y in range(top_px):
        alpha = round(top_alpha * (1 - y / top_px))
        draw.line((0, y, rgba.width, y), fill=(0, 0, 0, alpha))
    start = int(h * bottom_start)
    span = max(1, h - start)
    for y in range(start, h):
        alpha = round(bottom_alpha * ((y - start) / span))
        draw.line((0, y, rgba.width, y), fill=(0, 0, 0, alpha))
    return Image.alpha_composite(rgba, overlay).convert("RGB")


def prepare_background(source: Path, target: Path, kind: str, cta_alpha: int) -> None:
    target.parent.mkdir(parents=True, exist_ok=True)
    fitted = target.with_name(target.stem + "-fit.jpg")
    fit_cover(source, fitted)
    base = Image.open(fitted).convert("RGB")
    if kind == "cover":
        result = gradient_overlay(base, top_end=0.25, bottom_start=0.50, top_alpha=175, bottom_alpha=215)
    elif kind == "cta":
        overlay = Image.new("RGBA", base.size, (0, 0, 0, cta_alpha))
        result = Image.alpha_composite(base.convert("RGBA"), overlay).convert("RGB")
    else:
        result = gradient_overlay(base, top_end=0.25, bottom_start=0.3736, top_alpha=165, bottom_alpha=225)
    target.parent.mkdir(parents=True, exist_ok=True)
    result.save(target, quality=95)
    fitted.unlink(missing_ok=True)


def add_textbox(slide, left: float, top: float, width: float, height: float, *,
                text: str = "", font: str, size: float, bold: bool = False,
                align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, line_spacing: float | None = None):
    shape = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    frame = shape.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = anchor
    paragraphs = text.split("\n") or [""]
    for idx, value in enumerate(paragraphs):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.text = value
        p.alignment = align
        if line_spacing is not None:
            p.line_spacing = line_spacing
        if not p.runs:
            run = p.add_run()
            run.text = value
        for run in p.runs:
            run.font.name = font
            run.font.size = Pt(size)
            run.font.bold = bold
            run.font.color.rgb = WHITE
    return shape


def add_header(slide, font: str, category: str) -> None:
    add_textbox(slide, 0.3937, 0.3937, 2.4188, 0.2693, text="SRC Plus", font=font, size=16)
    add_textbox(slide, 2.8125, 0.3937, 2.4188, 0.2693, text=category, font=font, size=16, align=PP_ALIGN.RIGHT)


def add_source(slide, font: str, source: str) -> None:
    add_textbox(slide, 2.8125, 7.0811, 2.4188, 0.1683, text=source, font=font, size=10, align=PP_ALIGN.RIGHT)


def body_title_height(title: str, configured: float) -> tuple[float, float]:
    lines = title.count("\n") + 1
    if lines > 1 or len(title) > 20:
        return 4.4420, 0.9425
    return configured, 0.4713


def add_cta_text(slide, page: dict, design: dict, font: str) -> None:
    box = slide.shapes.add_textbox(
        Inches(design["cta"]["left"]), Inches(design["cta"]["top"]),
        Inches(design["cta"]["width"]), Inches(design["cta"]["height"]),
    )
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = frame.margin_right = frame.margin_top = frame.margin_bottom = 0
    frame.vertical_anchor = MSO_ANCHOR.TOP
    subject = page["title"]
    lines = [
        (subject, "에 대한"),
        ("", "더 자세한 이야기와, 다른 다양한 주제에 대한 리포트는"),
        ("SRC_Plus", " 페이지( https://srcplus.vercel.app/ )에서"),
        ("", "무료로 만나보실 수 있습니다."),
    ]
    for idx, (bold_text, normal_text) in enumerate(lines):
        p = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        p.alignment = PP_ALIGN.CENTER
        p.line_spacing = design["body"]["line_spacing"]
        if bold_text:
            run = p.add_run()
            run.text = bold_text
            run.font.name = font
            run.font.size = Pt(design["cta"]["font_size"])
            run.font.bold = True
            run.font.color.rgb = WHITE
        run = p.add_run()
        run.text = normal_text
        run.font.name = font
        run.font.size = Pt(design["cta"]["font_size"])
        run.font.bold = False
        run.font.color.rgb = WHITE


def build(manifest_path: Path, assets_dir: Path, design_path: Path, editorial_path: Path, output: Path) -> None:
    manifest = json.loads(manifest_path.read_text(encoding="utf-8"))
    design = json.loads(design_path.read_text(encoding="utf-8"))
    validate_manifest(manifest, load_editorial_rules(editorial_path))
    prs = Presentation()
    prs.slide_width = Inches(design["slide_width_in"])
    prs.slide_height = Inches(design["slide_height_in"])
    blank = prs.slide_layouts[6]
    font = design["font_family"]
    category = design["category_labels"].get(manifest.get("category"), str(manifest.get("category", "")).title())
    prepared = output.parent / "prepared-backgrounds"

    for page in manifest["pages"]:
        page_no = int(page["page_no"])
        kind = page["page_kind"]
        source = assets_dir / f"page-{page_no:02d}.png"
        bg = prepared / f"page-{page_no:02d}.jpg"
        prepare_background(source, bg, kind, int(design["cta"]["overlay_alpha"]))
        slide = prs.slides.add_slide(blank)
        slide.shapes.add_picture(str(bg), 0, 0, width=prs.slide_width, height=prs.slide_height)
        add_header(slide, font, category)
        add_source(slide, font, design["image_source"])

        if kind == "cover":
            add_textbox(
                slide, design["cover"]["title_left"], design["cover"]["title_top"],
                design["cover"]["title_width"], design["cover"]["title_height"],
                text=page["title"], font=font, size=design["cover"]["title_size"],
                bold=True, anchor=MSO_ANCHOR.BOTTOM,
            )
            add_textbox(
                slide, design["cover"]["title_left"], design["cover"]["subtitle_top"],
                design["cover"]["title_width"], design["cover"]["subtitle_height"],
                text=page["body"], font=font, size=16, anchor=MSO_ANCHOR.MIDDLE,
            )
        elif kind == "cta":
            add_cta_text(slide, page, design, font)
        else:
            title_top, title_height = body_title_height(page["title"], design["body"]["title_top"])
            add_textbox(
                slide, design["body"]["title_left"], title_top,
                design["body"]["title_width"], title_height,
                text=page["title"], font=font, size=design["body"]["title_size"],
                bold=True, anchor=MSO_ANCHOR.BOTTOM,
            )
            add_textbox(
                slide, design["body"]["title_left"], design["body"]["body_top"],
                design["body"]["title_width"], design["body"]["body_height"],
                text=page["body"], font=font, size=design["body"]["body_size"],
                line_spacing=design["body"]["line_spacing"],
            )

    output.parent.mkdir(parents=True, exist_ok=True)
    prs.save(output)
    print(f"built {output} with {len(manifest['pages'])} slides")


def main() -> None:
    parser = argparse.ArgumentParser()
    parser.add_argument("--manifest", required=True, type=Path)
    parser.add_argument("--assets", required=True, type=Path)
    parser.add_argument("--design", required=True, type=Path)
    parser.add_argument("--editorial", required=True, type=Path)
    parser.add_argument("--output", required=True, type=Path)
    args = parser.parse_args()
    build(args.manifest, args.assets, args.design, args.editorial, args.output)


if __name__ == "__main__":
    main()
